"""
Document Processing Pipeline
Handles various document formats and converts them to searchable chunks.
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
import re
from dataclasses import dataclass

# Document processing imports
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

@dataclass
class DocumentChunk:
    content: str
    metadata: Dict[str, Any]
    chunk_id: str

class DocumentProcessor:
    def __init__(self, 
                 chunk_size: int = 1000,
                 chunk_overlap: int = 200,
                 min_chunk_size: int = 100):
        """
        Initialize document processor.
        
        Args:
            chunk_size: Maximum size of each chunk
            chunk_overlap: Overlap between chunks
            min_chunk_size: Minimum size for a chunk to be kept
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size
        
        # Supported file extensions
        self.supported_extensions = {
            '.txt': self._process_text,
            '.md': self._process_text,
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.html': self._process_html,
            '.htm': self._process_html
        }
    
    def process_directory(self, directory_path: str) -> List[DocumentChunk]:
        """
        Process all supported documents in a directory.
        
        Args:
            directory_path: Path to directory containing documents
            
        Returns:
            List of document chunks
        """
        directory = Path(directory_path)
        all_chunks = []
        
        print(f"🔄 Processing documents in: {directory}")
        
        for file_path in directory.rglob("*"):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_extensions:
                try:
                    chunks = self.process_file(str(file_path))
                    all_chunks.extend(chunks)
                    print(f"✅ Processed: {file_path.name} ({len(chunks)} chunks)")
                except Exception as e:
                    print(f"❌ Error processing {file_path.name}: {str(e)}")
        
        print(f"🎉 Total chunks created: {len(all_chunks)}")
        return all_chunks
    
    def process_file(self, file_path: str) -> List[DocumentChunk]:
        """
        Process a single file into chunks.
        
        Args:
            file_path: Path to the file
            
        Returns:
            List of document chunks
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {extension}")
        
        # Extract text using appropriate processor
        text = self.supported_extensions[extension](file_path)
        
        if not text or len(text.strip()) < self.min_chunk_size:
            return []
        
        # Create base metadata
        metadata = {
            'source': str(file_path),
            'filename': file_path.name,
            'file_type': extension,
            'file_size': file_path.stat().st_size,
        }
        
        # Split text into chunks
        chunks = self._create_chunks(text, metadata)
        
        return chunks
    
    def _process_text(self, file_path: Path) -> str:
        """Process plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()
    
    def _process_pdf(self, file_path: Path) -> str:
        """Process PDF files with improved error handling."""
        if PyPDF2 is None:
            raise ImportError("PyPDF2 is required for PDF processing. Install with: pip install PyPDF2")
        
        # Use general warnings instead of PyPDF2-specific warnings
        import warnings
        warnings.filterwarnings('ignore')
        
        text = ""
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                total_pages = len(pdf_reader.pages)
                print(f"Processing PDF: {file_path.name} ({total_pages} pages)")
                
                # Add timeout protection for each page
                for page_num in range(total_pages):
                    try:
                        page = pdf_reader.pages[page_num]
                        # Extract text with error handling
                        try:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                        except Exception as e:
                            print(f"Warning: Could not extract text from page {page_num}: {e}")
                            continue
                    except Exception as e:
                        print(f"Warning: Could not process page {page_num}: {e}")
                        continue
                    
                    # Add progress indication
                    if page_num % 5 == 0:  # Show progress every 5 pages
                        print(f"📄 Progress: {page_num + 1}/{total_pages} pages")
            
            if not text.strip():
                print(f"Warning: No text extracted from PDF: {file_path}")
                return ""
            
            print(f"✅ Successfully processed PDF: {file_path.name}")
            return text
        
        except Exception as e:
            print(f"Error processing PDF {file_path}: {e}")
            return ""
    
    def _process_docx(self, file_path: Path) -> str:
        """Process DOCX files."""
        if DocxDocument is None:
            raise ImportError("python-docx is required for DOCX processing. Install with: pip install python-docx")
        
        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        return text
    
    def _process_pptx(self, file_path: Path) -> str:
        """Process PPTX files."""
        if Presentation is None:
            raise ImportError("python-pptx is required for PPTX processing. Install with: pip install python-pptx")
        
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text
    
    def _process_html(self, file_path: Path) -> str:
        """Process HTML files."""
        if BeautifulSoup is None:
            raise ImportError("beautifulsoup4 is required for HTML processing. Install with: pip install beautifulsoup4")
        
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            return soup.get_text()
    
    def _create_chunks(self, text: str, base_metadata: Dict[str, Any]) -> List[DocumentChunk]:
        """
        Split text into overlapping chunks.
        
        Args:
            text: Full text to chunk
            base_metadata: Metadata to apply to all chunks
            
        Returns:
            List of document chunks
        """
        # Clean and normalize text
        text = self._clean_text(text)
        
        if len(text) < self.min_chunk_size:
            return []
        
        chunks = []
        start = 0
        chunk_num = 0
        
        while start < len(text):
            # Calculate end position
            end = start + self.chunk_size
            
            # If we're at the end, take everything
            if end >= len(text):
                chunk_text = text[start:]
            else:
                # Try to break at sentence boundary
                chunk_text = text[start:end]
                
                # Look for sentence endings near the boundary
                sentence_ends = [m.end() for m in re.finditer(r'[.!?]\s+', chunk_text)]
                if sentence_ends:
                    # Use the last sentence ending as the boundary
                    actual_end = start + sentence_ends[-1]
                    chunk_text = text[start:actual_end]
                else:
                    # Fall back to word boundary
                    words = chunk_text.split()
                    if len(words) > 1:
                        chunk_text = ' '.join(words[:-1])
            
            # Only add chunk if it's large enough
            if len(chunk_text.strip()) >= self.min_chunk_size:
                chunk_metadata = base_metadata.copy()
                chunk_metadata.update({
                    'chunk_id': chunk_num,
                    'chunk_size': len(chunk_text),
                    'start_index': start,
                    'end_index': start + len(chunk_text)
                })
                
                chunks.append(DocumentChunk(
                    content=chunk_text.strip(),
                    metadata=chunk_metadata,
                    chunk_id=f"{base_metadata['filename']}_chunk_{chunk_num}"
                ))
                
                chunk_num += 1
            
            # Move start position
            start = start + len(chunk_text) - self.chunk_overlap
            
            # Prevent infinite loop
            if start + self.chunk_size >= len(text):
                break
        
        return chunks
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text

if __name__ == "__main__":
    # Test the document processor
    processor = DocumentProcessor(chunk_size=500, chunk_overlap=100)
    
    # Create sample text file for testing
    test_text = """
    This is a sample document for testing the document processor.
    It contains multiple paragraphs and sentences to demonstrate
    how the chunking algorithm works.
    
    The processor should split this text into manageable chunks
    while preserving sentence boundaries when possible.
    
    This helps maintain context and readability in the resulting chunks.
    """
    
    # Save sample file
    os.makedirs("test_data", exist_ok=True)
    with open("test_data/sample.txt", "w") as f:
        f.write(test_text)
    
    # Process the file
    chunks = processor.process_file("test_data/sample.txt")
    
    print(f"Created {len(chunks)} chunks:")
    for i, chunk in enumerate(chunks):
        print(f"\nChunk {i+1}:")
        print(f"Content: {chunk.content[:100]}...")
        print(f"Metadata: {chunk.metadata}") 